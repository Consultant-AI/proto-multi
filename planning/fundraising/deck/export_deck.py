#!/usr/bin/env python3
"""
Export Proto pitch deck to PDF and PowerPoint formats.

Requirements:
    pip install playwright python-pptx Pillow

For playwright, also run:
    playwright install chromium

Usage:
    python export_deck.py              # Export both PDF and PPT
    python export_deck.py --pdf        # Export PDF only
    python export_deck.py --ppt        # Export PPT only
"""

import asyncio
import argparse
import os
from pathlib import Path

# Colors matching the deck theme - professional blue palette
COLORS = {
    'bg_primary': '#0a0a0a',
    'bg_secondary': '#161616',
    'accent': '#4a9eff',
    'cyan': '#6cb6ff',
    'yellow': '#e5a635',
    'text_primary': '#f5f5f5',
    'text_secondary': '#888888',
}

SCRIPT_DIR = Path(__file__).parent
HTML_FILE = SCRIPT_DIR / 'pitch_deck.html'
PDF_OUTPUT = SCRIPT_DIR / 'Proto_Pitch_Deck.pdf'
PPT_OUTPUT = SCRIPT_DIR / 'Proto_Pitch_Deck.pptx'


async def export_pdf():
    """Export the HTML presentation to PDF using Playwright."""
    try:
        from playwright.async_api import async_playwright
    except ImportError:
        print("Error: playwright not installed. Run: pip install playwright && playwright install chromium")
        return False

    print("Exporting to PDF...")

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page()

        # Load the HTML file
        await page.goto(f'file://{HTML_FILE.absolute()}')
        await page.wait_for_load_state('networkidle')

        # Get total slides
        total_slides = await page.evaluate('document.querySelectorAll(".slide").length')

        # Create a temporary HTML that shows all slides for PDF
        pdf_html = f'''
        <!DOCTYPE html>
        <html>
        <head>
            <style>
                @import url('https://fonts.googleapis.com/css2?family=JetBrains+Mono:wght@400;500;600;700&family=Space+Grotesk:wght@300;400;500;600;700&display=swap');
                @page {{ size: 1920px 1080px; margin: 0; }}
                body {{ margin: 0; padding: 0; }}
                .page {{
                    width: 1920px;
                    height: 1080px;
                    page-break-after: always;
                    position: relative;
                    overflow: hidden;
                }}
                .page:last-child {{ page-break-after: avoid; }}
            </style>
        </head>
        <body></body>
        </html>
        '''

        # Navigate through slides and capture each
        slides_html = []
        for i in range(total_slides):
            await page.evaluate(f'goToSlide({i})')
            await page.wait_for_timeout(100)  # Wait for animation

            # Get the slide content
            slide_html = await page.evaluate('''() => {
                const slide = document.querySelector('.slide.active');
                const styles = document.querySelector('style').innerHTML;
                return {
                    content: slide.outerHTML,
                    styles: styles
                };
            }''')
            slides_html.append(slide_html)

        # Create combined PDF HTML
        combined_html = f'''
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <style>
                {slides_html[0]['styles']}
                @page {{ size: 1920px 1080px; margin: 0; }}
                body {{ margin: 0; padding: 0; background: #000; }}
                .slide {{
                    position: relative !important;
                    display: flex !important;
                    opacity: 1 !important;
                    width: 1920px;
                    height: 1080px;
                    page-break-after: always;
                }}
                .slide:last-child {{ page-break-after: avoid; }}
                .nav, .slide-counter, .controls-hint {{ display: none !important; }}
            </style>
        </head>
        <body>
            <div class="presentation" style="width: 1920px; background-size: 60px 60px;">
        '''

        for slide in slides_html:
            # Make each slide visible
            content = slide['content'].replace('class="slide"', 'class="slide active"')
            combined_html += content

        combined_html += '</div></body></html>'

        # Create new page with combined content
        pdf_page = await browser.new_page()
        await pdf_page.set_viewport_size({'width': 1920, 'height': 1080})
        await pdf_page.set_content(combined_html)
        await pdf_page.wait_for_load_state('networkidle')
        await pdf_page.wait_for_timeout(1000)  # Wait for fonts to load

        # Generate PDF
        await pdf_page.pdf(
            path=str(PDF_OUTPUT),
            width='1920px',
            height='1080px',
            print_background=True,
            scale=1
        )

        await browser.close()
        print(f"PDF exported: {PDF_OUTPUT}")
        return True


def export_ppt():
    """Export the presentation to PowerPoint format."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.util import Emu
    except ImportError as e:
        print(f"Error: python-pptx not installed or import error: {e}")
        print("Run: pip install python-pptx")
        return False

    print("Exporting to PowerPoint...")

    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Blank layout
    blank_layout = prs.slide_layouts[6]

    from pptx.dml.color import RGBColor

    def hex_to_rgb(hex_color):
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )

    def add_text_box(slide, left, top, width, height, text, font_size=18,
                     font_name='Arial', bold=False, color='#ffffff', align='left'):
        from pptx.enum.text import PP_ALIGN
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.bold = bold
        p.font.color.rgb = hex_to_rgb(color)
        if align == 'center':
            p.alignment = PP_ALIGN.CENTER
        elif align == 'right':
            p.alignment = PP_ALIGN.RIGHT
        return shape

    def add_rectangle(slide, left, top, width, height, fill_color=None, border_color=None):
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        else:
            shape.fill.background()
        if border_color:
            shape.line.color.rgb = hex_to_rgb(border_color)
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(color)

    # ===== SLIDE 1: Title =====
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1, COLORS['bg_primary'])
    add_text_box(slide1, 0.5, 0.5, 3, 0.5, 'PRE-SEED', 14, 'Consolas', False, COLORS['accent'])
    add_text_box(slide1, 0.5, 2.5, 12, 1.5, 'PROTO', 96, 'Consolas', True, COLORS['text_primary'])
    add_text_box(slide1, 0.5, 4.5, 12, 0.5, 'THE AUTONOMOUS COMPANY FACTORY', 20, 'Consolas', False, COLORS['text_secondary'])
    add_text_box(slide1, 0.5, 6.5, 6, 0.5, 'Nir Feinstein / January 2026', 14, 'Consolas', False, '#444444')

    # ===== SLIDE 2: The Opportunity =====
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2, COLORS['bg_primary'])
    add_text_box(slide2, 0.5, 0.5, 3, 0.5, 'THE OPPORTUNITY', 14, 'Consolas', False, COLORS['accent'])
    add_text_box(slide2, 0.5, 1.2, 12, 1, 'Software is eating labor.', 48, 'Arial', True, COLORS['text_primary'])
    add_text_box(slide2, 0.5, 2.8, 4, 1, '$600B', 72, 'Consolas', True, COLORS['text_primary'])
    add_text_box(slide2, 0.5, 3.8, 4, 0.5, 'Software market', 18, 'Arial', False, COLORS['text_secondary'])
    add_text_box(slide2, 4.5, 3.2, 1, 0.5, 'vs', 18, 'Consolas', False, '#444444')
    add_text_box(slide2, 5.5, 2.8, 5, 1, '$90T', 72, 'Consolas', True, COLORS['accent'])
    add_text_box(slide2, 5.5, 3.8, 5, 0.5, 'Labor market', 18, 'Arial', False, COLORS['text_secondary'])
    add_rectangle(slide2, 0.5, 5, 12, 1, '#0d1f15', COLORS['accent'])
    add_text_box(slide2, 0.7, 5.2, 11.5, 0.8, "We're not competing with software. We're replacing labor.", 24, 'Arial', False, COLORS['text_primary'])

    # ===== SLIDE 3: The Prize =====
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3, COLORS['bg_primary'])
    add_text_box(slide3, 0.5, 0.5, 3, 0.5, 'THE PRIZE', 14, 'Consolas', False, COLORS['accent'])
    add_text_box(slide3, 0.5, 1.2, 12, 1, 'The compounding machine.', 48, 'Arial', True, COLORS['text_primary'])
    # Flow boxes
    flow_items = ['Business', 'Revenue', 'Data', 'Better System', 'More Businesses']
    for i, item in enumerate(flow_items):
        x = 0.5 + i * 2.5
        color = COLORS['accent'] if i == 4 else '#1a1a1a'
        text_color = COLORS['accent'] if i == 4 else COLORS['text_secondary']
        add_rectangle(slide3, x, 3, 2, 0.8, COLORS['bg_secondary'], color)
        add_text_box(slide3, x, 3.2, 2, 0.5, item, 14, 'Consolas', False, text_color, 'center')
        if i < 4:
            add_text_box(slide3, x + 2.1, 3.2, 0.3, 0.5, '→', 16, 'Consolas', False, '#333333')
    add_rectangle(slide3, 0.5, 5, 12, 1, '#0a0a0a', COLORS['accent'])
    add_text_box(slide3, 0.7, 5.2, 11.5, 0.8, 'Run one business autonomously → run hundreds.', 24, 'Arial', False, COLORS['text_primary'])

    # ===== SLIDE 4: What Proto Is =====
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4, COLORS['bg_primary'])
    add_text_box(slide4, 0.5, 0.5, 3, 0.5, 'THE PRODUCT', 14, 'Consolas', False, COLORS['accent'])
    add_text_box(slide4, 0.5, 1.2, 12, 1, 'What Proto is.', 48, 'Arial', True, COLORS['text_primary'])
    # Left column - Built today
    add_text_box(slide4, 0.5, 2.3, 6, 0.4, 'BUILT TODAY', 12, 'Consolas', False, COLORS['accent'])
    built_items = ['159 specialist agents organized like a company', 'Project planning system for complex work',
                   'Cross-platform Mac + cloud Linux', 'Real computer use: files, terminal, GUI, browser']
    for i, item in enumerate(built_items):
        add_text_box(slide4, 0.5, 2.9 + i * 0.7, 0.4, 0.5, f'0{i+1}', 12, 'Consolas', False, COLORS['accent'])
        add_text_box(slide4, 1.0, 2.9 + i * 0.7, 5.5, 0.6, item, 16, 'Arial', False, COLORS['text_secondary'])
    # Right column - On roadmap
    add_text_box(slide4, 7, 2.3, 6, 0.4, 'ON ROADMAP', 12, 'Consolas', False, COLORS['yellow'])
    roadmap_items = ['Multi-computer orchestration', 'Self-improvement systems', 'Hybrid human delegation']
    for i, item in enumerate(roadmap_items):
        add_text_box(slide4, 7, 2.9 + i * 0.7, 0.4, 0.5, '→', 12, 'Consolas', False, COLORS['yellow'])
        add_text_box(slide4, 7.5, 2.9 + i * 0.7, 5, 0.6, item, 16, 'Arial', False, COLORS['text_secondary'])
    add_rectangle(slide4, 7, 5.2, 5.5, 0.8, COLORS['bg_secondary'], '#1a1a1a')
    add_text_box(slide4, 7.2, 5.4, 5, 0.5, "This isn't a concept. It's built.", 16, 'Arial', False, COLORS['text_primary'])

    # ===== SLIDE 5: Architecture =====
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5, COLORS['bg_primary'])
    add_text_box(slide5, 0.5, 0.5, 3, 0.5, 'ARCHITECTURE', 14, 'Consolas', False, COLORS['accent'])
    add_text_box(slide5, 0.5, 1.2, 12, 1, 'System overview.', 48, 'Arial', True, COLORS['text_primary'])
    # Architecture boxes
    add_rectangle(slide5, 4.5, 2.3, 4, 0.7, '#0a0a0a', COLORS['accent'])
    add_text_box(slide5, 4.5, 2.45, 4, 0.5, 'GLOBAL CEO (Opus 4.5)', 14, 'Consolas', False, COLORS['accent'], 'center')
    # Computer pool
    for i, label in enumerate(['Computer 1 / Product A', 'Computer 2 / Product B', 'Computer N / Shared']):
        x = 0.5 + i * 4.2
        add_rectangle(slide5, x, 3.3, 3.8, 0.6, COLORS['bg_secondary'], '#1a1a1a')
        add_text_box(slide5, x, 3.4, 3.8, 0.5, label, 11, 'Consolas', False, COLORS['text_secondary'], 'center')
    # Agent layer
    add_rectangle(slide5, 0.5, 4.2, 12, 1.2, COLORS['bg_secondary'], '#1a1a1a')
    add_text_box(slide5, 0.7, 4.3, 3, 0.3, '159 AGENTS', 10, 'Consolas', False, '#444444')
    agent_levels = [('C-SUITE', COLORS['accent']), ('DIRECTORS', COLORS['cyan']), ('SPECIALISTS', '#666666')]
    for i, (level, color) in enumerate(agent_levels):
        x = 1 + i * 3
        add_rectangle(slide5, x, 4.7, 2.2, 0.5, '#111111', color)
        add_text_box(slide5, x, 4.8, 2.2, 0.4, level, 11, 'Consolas', False, color, 'center')
        if i < 2:
            add_text_box(slide5, x + 2.3, 4.8, 0.5, 0.4, '→', 14, 'Consolas', False, '#333333')
    add_text_box(slide5, 10.2, 4.8, 2, 0.4, 'PEER-TO-PEER', 9, 'Consolas', False, '#666666', 'center')
    # Infrastructure
    infra = [('MESSAGE BUS', COLORS['yellow']), ('KNOWLEDGE HUB', COLORS['cyan']), ('REVENUE ENGINE', COLORS['accent'])]
    for i, (label, color) in enumerate(infra):
        x = 0.5 + i * 4.2
        add_rectangle(slide5, x, 5.7, 3.8, 0.6, COLORS['bg_secondary'], color)
        add_text_box(slide5, x, 5.8, 3.8, 0.5, label, 12, 'Consolas', False, color, 'center')

    # ===== SLIDE 6: Demo =====
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6, COLORS['bg_primary'])
    add_text_box(slide6, 0.5, 0.5, 3, 0.5, 'DEMONSTRATION', 14, 'Consolas', False, COLORS['accent'])
    add_text_box(slide6, 0.5, 1.2, 12, 1, 'See it work.', 48, 'Arial', True, COLORS['text_primary'])
    add_rectangle(slide6, 1.5, 2.5, 10, 2, COLORS['bg_secondary'], '#1a1a1a')
    add_text_box(slide6, 1.5, 2.8, 10, 0.4, 'LIVE DEMO', 14, 'Consolas', False, '#444444', 'center')
    add_text_box(slide6, 1.5, 3.3, 10, 0.6, 'Complex task → Multi-agent coordination → Output', 24, 'Arial', False, COLORS['text_primary'], 'center')
    add_text_box(slide6, 1.5, 4, 10, 0.4, '90-second demonstration', 16, 'Arial', False, COLORS['text_secondary'], 'center')
    # Terminal
    add_rectangle(slide6, 1.5, 5, 10, 1.2, COLORS['bg_secondary'], '#1a1a1a')
    add_text_box(slide6, 1.7, 5.2, 0.3, 0.4, '$', 14, 'Consolas', False, COLORS['accent'])
    add_text_box(slide6, 2.0, 5.2, 9, 0.4, 'proto "Build authentication system with tests"', 14, 'Consolas', False, COLORS['text_secondary'])
    add_text_box(slide6, 1.7, 5.7, 0.3, 0.4, '→', 14, 'Consolas', False, COLORS['accent'])
    add_text_box(slide6, 2.0, 5.7, 9, 0.4, 'CEO delegating to Backend Dev, Security Engineer, QA...', 14, 'Consolas', False, COLORS['text_secondary'])

    # ===== SLIDE 7: Business Model =====
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7, COLORS['bg_primary'])
    add_text_box(slide7, 0.5, 0.5, 3, 0.5, 'BUSINESS MODEL', 14, 'Consolas', False, COLORS['accent'])
    add_text_box(slide7, 0.5, 1.2, 12, 1, 'We run businesses.', 48, 'Arial', True, COLORS['text_primary'])
    add_rectangle(slide7, 0.5, 2.3, 12, 0.8, '#0a0a0a', COLORS['accent'])
    add_text_box(slide7, 0.7, 2.45, 11.5, 0.6, 'Proto is not SaaS. We operate autonomous companies.', 22, 'Arial', False, COLORS['text_primary'])
    # Business types grid
    biz_types = [('SAAS', 'Build and operate software products'),
                 ('AGENCIES', 'Marketing, dev, design services'),
                 ('CONTENT', 'Media, publishing, courses'),
                 ('E-COMMERCE', 'Product sourcing and sales'),
                 ('FREELANCE', 'Professional services at scale'),
                 ('PHYSICAL + AI', 'AI hires humans when needed')]
    for i, (title, desc) in enumerate(biz_types):
        row = i // 3
        col = i % 3
        x = 0.5 + col * 4.2
        y = 3.5 + row * 1.4
        border = COLORS['accent'] if i == 5 else '#1a1a1a'
        bg = '#0a0a0a' if i != 5 else '#0d1f15'
        add_rectangle(slide7, x, y, 3.8, 1.2, COLORS['bg_secondary'], border)
        add_text_box(slide7, x + 0.2, y + 0.2, 3.4, 0.4, title, 12, 'Consolas', True, COLORS['accent'] if i == 5 else COLORS['accent'])
        add_text_box(slide7, x + 0.2, y + 0.6, 3.4, 0.5, desc, 13, 'Arial', False, COLORS['text_secondary'])

    # ===== SLIDE 8: Market Gap =====
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8, COLORS['bg_primary'])
    add_text_box(slide8, 0.5, 0.5, 3, 0.5, 'MARKET GAP', 14, 'Consolas', False, COLORS['accent'])
    add_text_box(slide8, 0.5, 1.2, 12, 1, 'Everyone builds silos.', 48, 'Arial', True, COLORS['text_primary'])
    # Left column
    add_text_box(slide8, 0.5, 2.3, 6, 0.4, 'CURRENT AGENTS', 12, 'Consolas', False, '#444444')
    add_text_box(slide8, 0.5, 2.8, 6, 0.5, '40% of enterprise apps have agents — all domain silos', 14, 'Arial', False, COLORS['text_secondary'])
    agents = [('Cursor, Claude Code', 'only code'), ('Apollo, Clay', 'only sales'),
              ('Harvey', 'only legal'), ('Torq', 'only security')]
    for i, (name, domain) in enumerate(agents):
        add_text_box(slide8, 0.5, 3.5 + i * 0.6, 3, 0.5, name, 14, 'Consolas', False, COLORS['text_primary'])
        add_text_box(slide8, 4, 3.5 + i * 0.6, 2, 0.5, domain, 14, 'Consolas', False, '#444444')
    # Right column
    add_text_box(slide8, 7, 2.3, 6, 0.4, "WHAT'S MISSING", 12, 'Consolas', False, COLORS['accent'])
    missing = [('Cross-domain coordination', "They don't talk to each other"),
               ('Self-improvement', 'Each task makes Proto better at ALL tasks'),
               ('Full autonomy goal', 'Run businesses, not just tasks')]
    for i, (title, desc) in enumerate(missing):
        add_text_box(slide8, 7, 3.0 + i * 1.1, 0.4, 0.4, '×', 14, 'Consolas', False, COLORS['accent'])
        add_text_box(slide8, 7.5, 3.0 + i * 1.1, 5, 0.5, title, 16, 'Arial', True, COLORS['text_primary'])
        add_text_box(slide8, 7.5, 3.5 + i * 1.1, 5, 0.5, desc, 14, 'Arial', False, '#444444')

    # ===== SLIDE 9: Team + Timing =====
    slide9 = prs.slides.add_slide(blank_layout)
    set_background(slide9, COLORS['bg_primary'])
    add_text_box(slide9, 0.5, 0.5, 3, 0.5, 'TEAM + TIMING', 14, 'Consolas', False, COLORS['accent'])
    add_text_box(slide9, 0.5, 1.2, 12, 1, 'Why now. Why me.', 48, 'Arial', True, COLORS['text_primary'])
    # Grid boxes
    add_rectangle(slide9, 0.5, 2.5, 5.8, 1.5, '#0a0a0a', COLORS['accent'])
    add_text_box(slide9, 0.7, 2.7, 5.4, 0.4, 'SOLO BUILDER', 12, 'Consolas', True, COLORS['accent'])
    add_text_box(slide9, 0.7, 3.2, 5.4, 0.8, 'Built entire stack: backend, frontend, agent orchestration, cross-platform infrastructure. All working.', 14, 'Arial', False, COLORS['text_secondary'])
    add_rectangle(slide9, 6.7, 2.5, 5.8, 1.5, '#0a0a0a', COLORS['accent'])
    add_text_box(slide9, 6.9, 2.7, 5.4, 0.4, 'FIRST PRINCIPLES', 12, 'Consolas', True, COLORS['accent'])
    add_text_box(slide9, 6.9, 3.2, 5.4, 0.8, 'Started with end goal (autonomous companies), broke it down, built each piece systematically.', 14, 'Arial', False, COLORS['text_secondary'])
    # Window box
    add_rectangle(slide9, 0.5, 4.5, 12, 1, '#0a0a0a', COLORS['accent'])
    add_text_box(slide9, 0.7, 4.7, 11.5, 0.7, 'The window: Models work. Computer use works. Massive resources flowing in. Gap in the market.', 18, 'Arial', False, COLORS['text_primary'])
    add_text_box(slide9, 0.5, 5.8, 12, 0.5, "First system with real self-improvement compounds. That's the moat.", 18, 'Arial', False, COLORS['text_secondary'])

    # ===== SLIDE 10: The Ask =====
    slide10 = prs.slides.add_slide(blank_layout)
    set_background(slide10, COLORS['bg_primary'])
    add_text_box(slide10, 0.5, 0.5, 3, 0.5, 'THE ASK', 14, 'Consolas', False, COLORS['accent'])
    add_text_box(slide10, 0.5, 1.2, 12, 1, 'Investment terms.', 48, 'Arial', True, COLORS['text_primary'])
    # Ask grid
    ask_items = [('RAISING', '$500K', 'Close A → $1.5M Close B'),
                 ('TERMS', '$5M', 'Post-Money SAFE'),
                 ('USE OF FUNDS', '12mo', 'Founder + 2 Engineers + Compute'),
                 ('MILESTONES', 'Self-improve → Mastery', 'Domain mastery, then autonomous')]
    for i, (label, value, detail) in enumerate(ask_items):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 6.2
        y = 2.5 + row * 1.8
        add_rectangle(slide10, x, y, 5.8, 1.6, COLORS['bg_secondary'], '#1a1a1a')
        add_text_box(slide10, x + 0.2, y + 0.2, 5.4, 0.3, label, 11, 'Consolas', False, '#444444')
        add_text_box(slide10, x + 0.2, y + 0.6, 5.4, 0.6, value, 36 if i < 3 else 18, 'Consolas', True, COLORS['text_primary'])
        add_text_box(slide10, x + 0.2, y + 1.2, 5.4, 0.3, detail, 13, 'Arial', False, COLORS['text_secondary'])
    # CTA
    add_rectangle(slide10, 0.5, 6, 12, 0.8, '#0a0a0a', COLORS['accent'])
    add_text_box(slide10, 0.7, 6.15, 11.5, 0.6, 'Looking for investors who see the opportunity in autonomous companies.', 18, 'Arial', False, COLORS['text_primary'])

    # ===== SLIDE 11: End =====
    slide11 = prs.slides.add_slide(blank_layout)
    set_background(slide11, COLORS['bg_primary'])
    add_text_box(slide11, 0.5, 0.5, 12, 0.5, 'CONTACT', 14, 'Consolas', False, COLORS['accent'], 'center')
    add_text_box(slide11, 0.5, 2.5, 12, 1.5, 'Questions?', 72, 'Consolas', True, COLORS['text_primary'], 'center')
    add_text_box(slide11, 0.5, 4.5, 12, 0.5, 'nir@proto.ai', 20, 'Consolas', False, COLORS['text_secondary'], 'center')
    add_text_box(slide11, 0.5, 6.5, 12, 0.5, 'PROTO / THE AUTONOMOUS COMPANY FACTORY', 14, 'Consolas', False, '#444444', 'center')

    # Save
    prs.save(str(PPT_OUTPUT))
    print(f"PowerPoint exported: {PPT_OUTPUT}")
    return True


def main():
    parser = argparse.ArgumentParser(description='Export Proto pitch deck to PDF and PowerPoint')
    parser.add_argument('--pdf', action='store_true', help='Export PDF only')
    parser.add_argument('--ppt', action='store_true', help='Export PowerPoint only')
    args = parser.parse_args()

    # If neither flag specified, export both
    export_both = not args.pdf and not args.ppt

    if args.pdf or export_both:
        asyncio.run(export_pdf())

    if args.ppt or export_both:
        export_ppt()

    print("\nDone! Files exported to:")
    print(f"  {SCRIPT_DIR}")


if __name__ == '__main__':
    main()
